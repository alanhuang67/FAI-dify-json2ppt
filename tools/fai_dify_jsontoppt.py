import json, uuid, tempfile
import logging
import sys
import ast
from pathlib import Path
from typing import Any, Generator
from io import BytesIO

from pptx import Presentation
from pptx.slide import Slide
from dify_plugin import Tool
from dify_plugin.entities.tool import ToolInvokeMessage
from dify_plugin.plugin import Plugin
from dify_plugin.config.config import DifyPluginEnv

# 配置日志
# log_dir = Path(__file__).parent.parent / "logs"
# log_dir.mkdir(exist_ok=True)
# log_file = log_dir / "json2ppt.log"

# 配置日志处理器，输出到标准错误流
stream_handler = logging.StreamHandler(sys.stderr)
stream_handler.setLevel(logging.INFO)
formatter = logging.Formatter('%(asctime)s - %(name)s - %(levelname)s - %(message)s')
stream_handler.setFormatter(formatter)

# 配置日志记录器
logger = logging.getLogger(__name__)
logger.setLevel(logging.INFO)
logger.addHandler(stream_handler)

# 移除之前的控制台输出，因为它现在与stream_handler重复
# for handler in logger.handlers:
#     if isinstance(handler, logging.StreamHandler):
#         logger.removeHandler(handler)

logger.info("="*50)
logger.info("插件启动")

def _choose_layout(prs: Presentation, name: str):
    logger.info(f"正在查找布局: {name}")
    for layout in prs.slide_layouts:
        if layout.name == name:
            logger.info(f"找到布局: {name}")
            return layout
    logger.error(f"未找到布局: {name}")
    raise ValueError(f"layout_not_found: {name}")

def _fill_placeholders(slide, elements: dict):
    logger.info(f"开始填充占位符，元素数量: {len(elements)}")
    for shape in slide.shapes:
        if shape.name in elements:
            val = elements[shape.name]
            # 处理不同类型的值
            if isinstance(val, list):
                # 如果是列表，将每个元素转换为字符串并用换行符连接
                text = "\n".join(str(item) for item in val)
            else:
                # 如果是字符串，直接使用
                text = str(val)
            
            # 处理项目符号
            if "•" in text:
                # 将项目符号替换为 PowerPoint 的自动项目符号
                lines = text.split("\n")
                processed_lines = []
                for line in lines:
                    if line.strip().startswith("•"):
                        # 移除项目符号，PowerPoint 会自动添加
                        processed_lines.append(line.strip()[1:].strip())
                    else:
                        processed_lines.append(line)
                text = "\n".join(processed_lines)
            
            # 设置文本
            shape.text = text
            logger.info(f"已填充占位符: {shape.name}")

def render_pptx(slides_json: str, template_name: str | None, file_name: str) -> str:
    logger.info("开始渲染PPT")
    logger.info(f"模板名称: {template_name}")
    logger.info(f"输出文件名: {file_name}")
    
    # 打印原始输入 JSON 的详细信息，帮助调试
    logger.info(f"接收到的 slides_json 原始输入 (repr 形式，前1000字符): {repr(slides_json[:1000])}...")
    logger.info(f"接收到的 slides_json 长度: {len(slides_json)}")

    try:
        # 直接解析 JSON
        data = json.loads(slides_json)
        logger.info(f"JSON解析成功，幻灯片数量: {len(data['slides'])}")
    except json.JSONDecodeError as e:
        logger.error(f"JSON解析失败: {str(e)}", exc_info=True)
        raise ValueError(f"[json2pptx error] 无法解析 JSON 数据，请确保 JSON 格式正确且特殊字符已正确转义。原始错误: {str(e)}") from e
    
    tpl_path = Path(template_name) if template_name else Path(__file__).parent.parent / "templates" / "default.pptx"
    logger.info(f"模板路径: {tpl_path}")
    
    if not tpl_path.is_file():
        logger.error(f"模板文件不存在: {tpl_path}")
        raise FileNotFoundError(f"template_not_found: {tpl_path}")
    
    try:
        prs = Presentation(str(tpl_path))
        logger.info("成功加载模板")
    except Exception as e:
        logger.error(f"加载模板失败: {str(e)}")
        raise
    
    for i, slide_spec in enumerate(data["slides"], 1):
        logger.info(f"正在处理第 {i} 张幻灯片")
        try:
            slide = prs.slides.add_slide(_choose_layout(prs, slide_spec["layout"]))
            _fill_placeholders(slide, slide_spec.get("elements", {}))
        except Exception as e:
            logger.error(f"处理第 {i} 张幻灯片时出错: {str(e)}")
            raise
    
    out_path = Path(tempfile.gettempdir()) / f"{uuid.uuid4()}_{file_name}"
    if out_path.suffix.lower() != ".pptx":
        out_path = out_path.with_suffix(".pptx")
    
    logger.info(f"正在保存PPT到: {out_path}")
    try:
        prs.save(out_path)
        logger.info("PPT保存成功")
    except Exception as e:
        logger.error(f"保存PPT失败: {str(e)}")
        raise
    
    return str(out_path)

class FaiDifyJsonTopptTool(Tool):
    def _invoke(self, params: dict[str, Any]) -> Generator[ToolInvokeMessage, None, None]:
        """处理工具调用"""
        try:
            logger.info("开始处理工具调用")
            logger.info(f"输入参数: {params}")
            
            # 获取 JSON 数据
            slides_json = params.get('slides_json', '')
            if not slides_json:
                logger.error("未提供 slides_json 参数")
                yield self.create_text_message("[json2pptx error] 未提供 slides_json 参数")
                return
            
            # 获取用户输入的档名，如果未提供则使用默认值
            output_filename = params.get('file_name', 'presentation')
            # 确保档名以 .pptx 结尾
            if not output_filename.lower().endswith(".pptx"):
                output_filename += ".pptx"

            # 获取模板名称
            template_name = params.get('template_name', None)
            logger.info(f"接收到的 template_name 参数: {template_name}")

            # 渲染 PPTX
            pptx_data = self.render_pptx(slides_json, template_name)
            logger.info("PPTX 渲染完成")
            
            # 先发送一个成功文本消息
            yield self.create_text_message("PowerPoint 演示文稿 \'" + output_filename + "\' 已成功生成！")

            # 创建并发送文件消息
            yield self.create_blob_message(
                blob=pptx_data,
                meta={
                    "mime_type": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
                    "filename": output_filename # 使用动态档名
                }
            )
            logger.info("文件消息创建完成")
            
        except Exception as e:
            logger.error(f"处理工具调用时出错: {str(e)}")
            yield self.create_text_message(f"[json2pptx error] {e}")
            return

    def _get_layout(self, prs: Presentation, layout_name: str):
        """根据布局名称获取布局"""
        for layout in prs.slide_layouts:
            if layout.name == layout_name:
                return layout
        return None

    def _fill_placeholders(self, slide: Slide, elements: dict):
        """填充页面占位符"""
        for placeholder_name, content in elements.items():
            for shape in slide.placeholders:
                if shape.name == placeholder_name:
                    if isinstance(content, list):
                        # 处理列表内容
                        text_frame = shape.text_frame
                        text_frame.clear()
                        for item in content:
                            p = text_frame.add_paragraph()
                            p.text = str(item)
                            p.level = 0  # 默认级别
                    else:
                        # 处理字符串内容
                        shape.text = str(content)
                    break

    def render_pptx(self, slides_json: str, template_name: str | None) -> bytes:
        """渲染 PPTX 文件"""
        try:
            logger.info("开始渲染 PPTX")
            logger.info(f"输入 JSON: {slides_json}")
            logger.info(f"接收到的 template_name: {template_name}")
            
            # 解析 JSON
            data = json.loads(slides_json)
            logger.info(f"解析后的数据: {data}")
            
            # 确定模板路径
            base_template_dir = Path(__file__).parent.parent / "templates"
            tpl_path = None
            if template_name:
                # 确保自定义模板名称包含 .pptx 后缀
                if not template_name.lower().endswith(".pptx"):
                    template_name += ".pptx"
                tpl_path = base_template_dir / template_name
                logger.info(f"尝试加载自定义模板: {tpl_path}")
            else:
                tpl_path = base_template_dir / "default.pptx"
                logger.info(f"未指定模板，加载默认模板: {tpl_path}")

            if not tpl_path.is_file():
                logger.error(f"模板文件不存在: {tpl_path}")
                raise FileNotFoundError(f"template_not_found: {tpl_path}")

            # 创建演示文稿 (从模板加载)
            prs = Presentation(str(tpl_path))
            logger.info("成功加载演示文稿模板")
            
            # 处理每一页
            for slide_data in data.get('slides', []):
                logger.info(f"处理页面: {slide_data}")
                layout_name = slide_data.get('layout', 'Title and Content')
                elements = slide_data.get('elements', {})
                
                # 获取布局
                layout = self._get_layout(prs, layout_name)
                if not layout:
                    logger.warning(f"未找到布局 {layout_name}，使用默认布局")
                    layout = prs.slide_layouts[1]  # 默认使用 Title and Content
                
                # 创建页面
                slide = prs.slides.add_slide(layout)
                logger.info(f"创建页面，使用布局: {layout_name}")
                
                # 填充内容
                self._fill_placeholders(slide, elements)
                logger.info("页面内容填充完成")
            
            # 保存到内存
            pptx_stream = BytesIO()
            prs.save(pptx_stream)
            pptx_stream.seek(0)
            pptx_data = pptx_stream.getvalue()
            logger.info(f"PPTX 文件生成完成，大小: {len(pptx_data)} 字节")
            
            return pptx_data
            
        except Exception as e:
            logger.error(f"渲染 PPTX 时出错: {str(e)}")
            raise
